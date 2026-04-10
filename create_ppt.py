#!/usr/bin/env python3
"""Generate presentation PPT for Prof. Hait meeting (Apr 10, 2026)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import csv
import os

BASE = os.path.dirname(os.path.abspath(__file__))
RESULTS = os.path.join(BASE, "journal_experiments", "results")
FIGURES = os.path.join(BASE, "figures")

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x00, 0x80, 0x00)
RED = RGBColor(0xCC, 0x00, 0x00)
HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
ROW_ALT = RGBColor(0xE8, 0xF0, 0xFE)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=None):
    """Dark blue title bar at top."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p2.alignment = PP_ALIGN.LEFT


def add_text_box(slide, left, top, width, height, text, size=16, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, items, left=Inches(0.6), top=Inches(1.4), width=Inches(12), size=18, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = spacing
        p.level = 0
    return tf


def add_table(slide, headers, rows, left=Inches(0.4), top=Inches(1.5), width=Inches(12.5), row_height=Inches(0.4)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows)
    table = table_shape.table

    # Set column widths evenly
    col_w = int(width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.CENTER
    return table


def fmt(val, pct=True):
    """Format float value as percentage."""
    try:
        v = float(val)
        if pct:
            return f"{v*100:.2f}%"
        return f"{v:.4f}"
    except:
        return str(val)


def read_csv(name):
    path = os.path.join(RESULTS, name)
    with open(path) as f:
        return list(csv.DictReader(f))


# ========== SLIDE 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5), prs.slide_width, Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

add_text_box(slide, Inches(0.5), Inches(2.0), Inches(12), Inches(1),
             "A Grouping-Based Ensemble Method with\nNSGA-II for Lung Histopathology Classification",
             size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(3.5), Inches(12), Inches(0.8),
             "Anshul Saxena  |  Advik Kashi Vishwanath  |  Kritii Gupta  |  Chandrima Debnath",
             size=16, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(4.1), Inches(12), Inches(0.5),
             "Supervisor: Prof. Swati Hait  |  BITS Pilani, Hyderabad Campus",
             size=14, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
             "Progress Meeting — April 10, 2026",
             size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Skip BITS logo (file may be corrupt)


# ========== SLIDE 2: Agenda ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Agenda")
add_bullet_slide(slide, [
    "1.  Problem Statement & Dataset Overview",
    "2.  Proposed Architecture (Modified Algorithm)",
    "3.  Key Modifications — Grouping Operator, Adaptive GA, NSGA-II, GP Ensemble",
    "4.  Experiment Results (Tables 2–15)",
    "5.  Analysis of Prof. Hait's Observations",
    "6.  Table 16 Status & Next Steps",
], size=20)


# ========== SLIDE 3: Problem & Dataset ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Problem Statement & Dataset")
add_bullet_slide(slide, [
    "Task: 3-class lung histopathology classification (ACA / Normal / SCC)",
    "Dataset: LC25000 — 15,000 images (5,000 per class), 224×224 RGB",
    "Split: 70% Train / 10% Validation / 20% Test (stratified)",
    "",
    "Challenges:",
    "    • High-dimensional feature space from CNN backbones (512–2048 dims)",
    "    • Need effective feature selection to remove redundancy",
    "    • Combine multiple classifiers for robust predictions",
], size=18)


# ========== SLIDE 4: Architecture ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Proposed Architecture — Modified Algorithm")

arch_path = os.path.join(FIGURES, "methodology_architecture.png")
if os.path.exists(arch_path):
    slide.shapes.add_picture(arch_path, Inches(0.5), Inches(1.3), width=Inches(12.3))
else:
    add_bullet_slide(slide, [
        "CNN Backbone (4) → Attention (SE) → Feature Extraction",
        "→ Feature Selection (Baseline GA / Adaptive GA / NSGA-II)",
        "→ Classification (KNN, SVM, RF, LR, XGBoost)",
        "→ Ensemble Fusion (Majority Vote, Weighted, GP1–GP4)",
    ])


# ========== SLIDE 5: Key Modifications ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Key Modifications (New Functions)")

items = [
    "1. Grouping Operator for GA & NSGA-II",
    "     Features split into 4 quartile groups: [0,D/4), [D/4,D/2), [D/2,3D/4), [3D/4,D)",
    "     Crossover & mutation respect group boundaries → preserves spatial/channel structure",
    "",
    "2. Adaptive GA (rates adjust based on population diversity)",
    "     High diversity → exploit (high crossover, low mutation)",
    "     Low diversity → explore (low crossover, high mutation)",
    "",
    "3. NSGA-II Multi-Objective Feature Selection",
    "     Objective 1: Maximize classification accuracy",
    "     Objective 2: Minimize number of selected features",
    "     Uses Pareto-optimal front + crowding distance",
    "",
    "4. GP Ensemble Fusion (GP1–GP4 + Weighted Fusion)",
    "     Epsilon-weighted classifier probabilities for combining 5 classifiers",
]
add_bullet_slide(slide, items, size=16, spacing=Pt(4))


# ========== SLIDE 6: GP Operators Detail ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "GP Ensemble Fusion Operators")

items = [
    "5 Classifiers: KNN (k=5), SVM (RBF), Random Forest, Logistic Regression, XGBoost",
    "",
    "Epsilon Weights: Classifiers ranked by validation accuracy,",
    "     T_j = T_{j-1} × acc_{j-1}, then normalized: ε = T / Σ(T)",
    "",
    "GP1 (Maximum):       1 − min(1−z₁, 1−z₂, 1−z₃, 1−z₄)",
    "GP2 (Algebraic):     1 − (1−z₁)(1−z₂)(1−z₃)(1−z₄)",
    "GP3 (Ratio):           1 − (1 + z₁z₂z₃z₄) / ((1−z₁)(1−z₂)(1−z₃)(1−z₄))",
    "GP4 (Weighted Sum): (z₁+z₂+z₃+z₄) / (1 + z₁z₂z₃z₄)",
    "",
    "Majority Vote: Mode across all 5 classifiers",
    "Weighted Fusion: P_ens(c) = Σ(w_i × P_i(c)) / (1 + Π(w_i × P_i(c)))",
]
add_bullet_slide(slide, items, size=16, spacing=Pt(4))


# ========== SLIDE 7: Table 2 — Classifier Performance ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Table 2 — Individual Classifier Performance",
              "DenseNet121 + SE Attention + Adaptive GA (KNN, SVM, RF)")

data = read_csv("table_2.csv")
classifiers = ["KNN", "SVM", "RF"]
headers = ["Classifier", "Accuracy", "Precision", "Recall", "F1", "Specificity"]
rows = []
for i, row in enumerate(data):
    if i >= 3:
        break
    rows.append([
        classifiers[i],
        fmt(row['accuracy']),
        fmt(row['precision']),
        fmt(row['recall']),
        fmt(row['f1']),
        fmt(row['specificity']),
    ])
add_table(slide, headers, rows, top=Inches(1.5))


# ========== SLIDE 8: Table 5 — Attention Comparison ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Table 5 — Attention Mechanism Comparison",
              "DenseNet121 backbone + Adaptive GA + KNN")

data = read_csv("table_5.csv")
headers = ["Attention", "Accuracy", "Precision", "Recall", "F1", "Specificity", "Params"]
rows = []
for row in data:
    rows.append([
        row.get('Attention', row.get('attention', '')),
        fmt(row['Accuracy'] if 'Accuracy' in row else row['accuracy']),
        fmt(row['Precision'] if 'Precision' in row else row['precision']),
        fmt(row['Recall'] if 'Recall' in row else row['recall']),
        fmt(row['F1'] if 'F1' in row else row['f1']),
        fmt(row['Specificity'] if 'Specificity' in row else row['specificity']),
        row.get('Params', row.get('params', 'N/A')),
    ])
add_table(slide, headers, rows, top=Inches(1.5))


# ========== SLIDE 9: Table 6 — GA vs NSGA-II ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Table 6 — GA vs NSGA-II Feature Selection")

data = read_csv("table_6.csv")
headers = ["Method", "Accuracy", "Precision", "Recall", "F1", "Specificity", "Features", "Reduction"]
rows = []
for row in data:
    method = row.get('Method', row.get('method', ''))
    feats = row.get('Features', row.get('features', ''))
    reduct = row.get('Reduction', row.get('reduction', ''))
    rows.append([
        method,
        fmt(row.get('Accuracy', row.get('accuracy', ''))),
        fmt(row.get('Precision', row.get('precision', ''))),
        fmt(row.get('Recall', row.get('recall', ''))),
        fmt(row.get('F1', row.get('f1', ''))),
        fmt(row.get('Specificity', row.get('specificity', ''))),
        feats,
        reduct,
    ])
add_table(slide, headers, rows, top=Inches(1.5))


# ========== SLIDE 10: Table 7 & 8 — Grouping Operator Effect ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Tables 7 & 8 — Grouping Operator & GA Variants")

# Table 8
data8 = read_csv("table_8.csv")
headers = ["GA Variant", "Accuracy", "Precision", "Recall", "F1", "Features"]
rows = []
for row in data8:
    rows.append([
        row.get('Method', row.get('method', row.get('Variant', ''))),
        fmt(row.get('Accuracy', row.get('accuracy', ''))),
        fmt(row.get('Precision', row.get('precision', ''))),
        fmt(row.get('Recall', row.get('recall', ''))),
        fmt(row.get('F1', row.get('f1', ''))),
        row.get('Features', row.get('features', '')),
    ])
add_table(slide, headers, rows, top=Inches(1.5))


# ========== SLIDE 11: Table 9 & 10 — NSGA-II + Comprehensive FS ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Table 10 — Comprehensive Feature Selection Comparison",
              "Baseline GA vs Adaptive GA vs NSGA-II")

data = read_csv("table_10.csv")
headers = ["Method", "Accuracy", "Features", "Reduction", "Time (s)"]
rows = []
for row in data:
    rows.append([
        row.get('Method', row.get('method', '')),
        fmt(row.get('Accuracy', row.get('accuracy', ''))),
        row.get('Features', row.get('features', '')),
        row.get('Reduction', row.get('reduction', '')),
        row.get('Time (s)', row.get('time', '')),
    ])
add_table(slide, headers, rows, top=Inches(1.5))

add_bullet_slide(slide, [
    "Key Insight: NSGA-II achieves 69.7% feature reduction (best) but takes longest",
    "Baseline GA: 50.6% reduction, slightly higher accuracy (99.73%)",
    "Adaptive GA: 50.2% reduction, fastest convergence (1186s)",
], left=Inches(0.6), top=Inches(3.5), size=16)


# ========== SLIDE 12: Table 11-13 — Sensitivity Analysis ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Tables 11–13 — Hyperparameter Sensitivity Analysis")

# Pop size (Table 11)
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(4), Inches(0.4),
             "Table 11: Population Size", size=14, bold=True, color=ACCENT_BLUE)
data11 = read_csv("table_11.csv")
h11 = ["Pop", "Accuracy", "Features"]
r11 = [[row.get('Pop Size', row.get('pop_size', '')),
        fmt(row.get('Accuracy', row.get('accuracy', ''))),
        row.get('Features', row.get('features', ''))] for row in data11]
add_table(slide, h11, r11, left=Inches(0.3), top=Inches(1.7), width=Inches(3.8), row_height=Inches(0.35))

# Generations (Table 12)
add_text_box(slide, Inches(4.5), Inches(1.3), Inches(4), Inches(0.4),
             "Table 12: Generations", size=14, bold=True, color=ACCENT_BLUE)
data12 = read_csv("table_12.csv")
h12 = ["Gens", "Accuracy", "Features"]
r12 = [[row.get('Generations', row.get('generations', '')),
        fmt(row.get('Accuracy', row.get('accuracy', ''))),
        row.get('Features', row.get('features', ''))] for row in data12]
add_table(slide, h12, r12, left=Inches(4.3), top=Inches(1.7), width=Inches(3.8), row_height=Inches(0.35))

# KNN k (Table 13)
add_text_box(slide, Inches(8.7), Inches(1.3), Inches(4), Inches(0.4),
             "Table 13: KNN k Value", size=14, bold=True, color=ACCENT_BLUE)
data13 = read_csv("table_13.csv")
h13 = ["k", "Accuracy", "F1"]
r13 = [[row.get('k', ''),
        fmt(row.get('Accuracy', row.get('accuracy', ''))),
        fmt(row.get('F1', row.get('f1', '')))] for row in data13]
add_table(slide, h13, r13, left=Inches(8.5), top=Inches(1.7), width=Inches(3.8), row_height=Inches(0.35))

# Summary
add_bullet_slide(slide, [
    "Population size: Stable 99.7–99.87%. Pop=20 or 60 slightly better",
    "Generations: Converges by 50 gens. No improvement at 100",
    "KNN k: k=5 and k=7 optimal. k=3,9 slightly lower (99.6%)",
], left=Inches(0.6), top=Inches(4.5), size=16)


# ========== SLIDE 13: Table 14 — Backbone Comparison ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Table 14 — CNN Backbone Comparison")

data = read_csv("table_14.csv")
headers = ["Backbone", "Feature Dim", "Accuracy", "F1", "Specificity", "Selected Features"]
rows = []
for row in data:
    rows.append([
        row.get('Backbone', row.get('backbone', '')),
        row.get('Feature Dim', row.get('feature_dim', '')),
        fmt(row.get('Accuracy', row.get('accuracy', ''))),
        fmt(row.get('F1', row.get('f1', ''))),
        fmt(row.get('Specificity', row.get('specificity', ''))),
        row.get('Selected', row.get('selected', '')),
    ])
add_table(slide, headers, rows, top=Inches(1.5))

add_bullet_slide(slide, [
    "EfficientNetB0: Best accuracy (99.90%) with 1280-dim features",
    "DenseNet121: Strong (99.70%) with compact 1024-dim",
    "VGG16: Lowest (96.77%) — older architecture, fewer features",
], left=Inches(0.6), top=Inches(3.5), size=16)


# ========== SLIDE 14: Table 15 — Ensemble Comparison ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Table 15 — Ensemble Fusion Strategies",
              "5 Classifiers: KNN, SVM, RF, LR, XGBoost + 6 Ensemble Methods")

data = read_csv("table_15.csv")
headers = ["Method", "Accuracy", "Precision", "Recall", "F1", "Specificity"]
rows = []
for row in data:
    method = row.get('Method', row.get('method', row.get('Ensemble', '')))
    rows.append([
        method,
        fmt(row.get('Accuracy', row.get('accuracy', ''))),
        fmt(row.get('Precision', row.get('precision', ''))),
        fmt(row.get('Recall', row.get('recall', ''))),
        fmt(row.get('F1', row.get('f1', ''))),
        fmt(row.get('Specificity', row.get('specificity', ''))),
    ])
add_table(slide, headers, rows, top=Inches(1.5))


# ========== SLIDE 15: NSGA-II Convergence ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "NSGA-II Convergence & Pareto Front")

conv_path = os.path.join(FIGURES, "nsga2_convergence_combined.png")
pareto_path = os.path.join(FIGURES, "nsga2_pareto_evolution.png")

if os.path.exists(conv_path):
    slide.shapes.add_picture(conv_path, Inches(0.3), Inches(1.3), width=Inches(6.3))
if os.path.exists(pareto_path):
    slide.shapes.add_picture(pareto_path, Inches(6.8), Inches(1.3), width=Inches(6.3))


# ========== SLIDE 16: Confusion Matrix ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Confusion Matrix & Parameter Sensitivity")

cm_path = os.path.join(FIGURES, "confusion_matrix.png")
sens_path = os.path.join(FIGURES, "param_sensitivity_combined.png")

if os.path.exists(cm_path):
    slide.shapes.add_picture(cm_path, Inches(0.3), Inches(1.3), width=Inches(6.0))
if os.path.exists(sens_path):
    slide.shapes.add_picture(sens_path, Inches(6.5), Inches(1.3), width=Inches(6.5))


# ========== SLIDE 17: 10-Fold CV (Table 4) ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Table 4 — 10-Fold Cross-Validation",
              "DenseNet121 + SE + Adaptive GA + KNN")

data = read_csv("table_4.csv")
headers = ["Fold", "Accuracy", "Precision", "Recall", "F1", "Specificity"]
rows = []
for row in data:
    fold = row.get('Fold', row.get('fold', ''))
    rows.append([
        fold,
        fmt(row.get('Accuracy', row.get('accuracy', ''))),
        fmt(row.get('Precision', row.get('precision', ''))),
        fmt(row.get('Recall', row.get('recall', ''))),
        fmt(row.get('F1', row.get('f1', ''))),
        fmt(row.get('Specificity', row.get('specificity', ''))),
    ])
add_table(slide, headers, rows, top=Inches(1.4), row_height=Inches(0.35))


# ========== SLIDE 18: Prof's Concerns ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Addressing Prof. Hait's Observations (March 31)")

items = [
    "Q1: \"Have you applied the new grouping function?\"",
    "   ✓ Yes — Grouping operator implemented in both Adaptive GA and NSGA-II",
    "   ✓ Features split into 4 quartile groups; crossover/mutation respect group boundaries",
    "",
    "Q2: \"Most values in many tables are exactly the same\"",
    "   → Accuracy values are very close (99.60–99.87%) because LC25000 dataset is relatively easy",
    "   → Small differences visible in features selected and per-class metrics",
    "   → Need to investigate: are features actually changing between GA variants?",
    "",
    "Q3: \"NSGA-II and Simple GA have not affected results — not even a small change\"",
    "   → Feature REDUCTION differs significantly: Baseline GA 50.6% vs NSGA-II 69.7%",
    "   → But classification accuracy stays similar (~99.7%) due to dataset saturation",
    "   → Key benefit of NSGA-II: fewer features, same performance = more efficient model",
    "",
    "Q4: \"Have new grouping operators improved results?\"",
    "   → Grouping maintains accuracy while providing better feature diversity",
    "   → Table 7: With grouping 99.77% vs without 99.70% (marginal improvement)",
    "   → Main value: structural coherence in selected features, not raw accuracy gain",
]
add_bullet_slide(slide, items, size=14, spacing=Pt(2))


# ========== SLIDE 19: Table 16 Status ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Table 16 — Full Pipeline Status")

items = [
    "Table 16: 4 Backbones × SE Attention × NSGA-II × 5 Classifiers × GP Ensemble",
    "",
    "Status: NOT COMPLETED — Out of Memory (OOM)",
    "",
    "Hardware used: NVIDIA L4 GPU, 32 GB RAM, 8-core CPU",
    "Issue: NSGA-II on concatenated 4-backbone features (1024+2048+512+1280 = 4864 dims)",
    "       exceeds available memory even with reduced pop/gen",
    "",
    "Options to resolve:",
    "   1. Access machine with ≥64 GB RAM",
    "   2. Process backbones sequentially instead of concatenating",
    "   3. Apply dimensionality reduction (PCA) before NSGA-II",
    "   4. Use smaller population size with batched fitness evaluation",
]
add_bullet_slide(slide, items, size=17, spacing=Pt(4))


# ========== SLIDE 20: Summary & Next Steps ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Summary & Next Steps")

items = [
    "Completed:",
    "   ✓ Tables 2–15 all generated with new functions (grouping, adaptive GA, NSGA-II, GP ensemble)",
    "   ✓ 5 classifiers (KNN, SVM, RF, LR, XGBoost) + 6 ensemble methods",
    "   ✓ Comprehensive ablation studies (attention, backbone, hyperparameters)",
    "   ✓ Best accuracy: 99.90% (EfficientNetB0), 99.70% (DenseNet121 default)",
    "",
    "Pending:",
    "   ⏳ Table 16 — needs higher RAM machine",
    "   ⏳ Verify feature selection actually changes features across GA variants",
    "   ⏳ Add old vs new function comparison tables",
    "",
    "Next Steps:",
    "   → Resolve Table 16 OOM issue",
    "   → Finalize results in Overleaf",
    "   → Complete paper draft for submission",
]
add_bullet_slide(slide, items, size=17, spacing=Pt(4))


# ========== SLIDE 21: Thank You ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

add_text_box(slide, Inches(0.5), Inches(2.8), Inches(12), Inches(1),
             "Thank You!", size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.5), Inches(3.7), Inches(12), Inches(0.5),
             "Questions & Discussion", size=20, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)


# Save
out_path = os.path.join(BASE, "presentation_apr10.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
